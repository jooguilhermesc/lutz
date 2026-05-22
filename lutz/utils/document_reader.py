"""Multi-format text extraction for context files.

Supports PDF, DOCX, XLSX, and PPTX. Returns a list of (page_num, text) tuples
compatible with the PDFProcessor page format so chunking pipelines can reuse
the same infrastructure.
"""
from __future__ import annotations

from pathlib import Path


def extract_pages(path: Path) -> list[tuple[int, str]]:
    """Extract text from a file and return a list of (page_num, text) tuples.

    Supported formats:
      .pdf  — via pymupdf (primary) or pdfplumber fallback
      .docx — via python-docx (paragraphs grouped into virtual pages)
      .xlsx — via openpyxl (each sheet becomes a virtual page)
      .pptx — via python-pptx (each slide becomes a virtual page)

    Raises:
        ValueError: if the file extension is not supported.
        ImportError: if the required library for the format is not installed.
    """
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return _extract_pdf(path)
    elif suffix == ".docx":
        return _extract_docx(path)
    elif suffix in (".xlsx", ".xls"):
        return _extract_xlsx(path)
    elif suffix == ".pptx":
        return _extract_pptx(path)
    elif suffix in (".txt", ".md", ".markdown", ".rst", ".csv"):
        return _extract_text(path)
    else:
        raise ValueError(
            f"Unsupported file type: '{suffix}'. "
            "Supported: .pdf, .docx, .xlsx, .xls, .pptx, .txt, .md, .csv"
        )


def _extract_pdf(path: Path) -> list[tuple[int, str]]:
    """Extract pages from a PDF using pymupdf, falling back to pdfplumber."""
    try:
        import fitz  # pymupdf

        doc = fitz.open(str(path))
        pages = []
        for i, page in enumerate(doc, 1):
            text = page.get_text("text").strip()
            if text:
                pages.append((i, text))
        doc.close()
        return pages
    except ImportError:
        pass

    try:
        import pdfplumber

        pages = []
        with pdfplumber.open(str(path)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = (page.extract_text() or "").strip()
                if text:
                    pages.append((i, text))
        return pages
    except ImportError:
        pass

    raise ImportError(
        "PDF extraction requires pymupdf or pdfplumber. "
        "Install with: pip install pymupdf"
    )


def _extract_docx(path: Path) -> list[tuple[int, str]]:
    """Extract text from a DOCX file, grouping paragraphs into virtual pages."""
    try:
        from docx import Document
    except ImportError:
        raise ImportError(
            "DOCX extraction requires python-docx. "
            "Install with: pip install python-docx"
        )

    doc = Document(str(path))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    # Group into virtual pages of ~40 paragraphs each
    _PAGE_SIZE = 40
    pages = []
    for i in range(0, max(len(paragraphs), 1), _PAGE_SIZE):
        chunk = paragraphs[i : i + _PAGE_SIZE]
        if chunk:
            pages.append((len(pages) + 1, "\n\n".join(chunk)))
    return pages


def _extract_xlsx(path: Path) -> list[tuple[int, str]]:
    """Extract text from an XLSX file, one virtual page per sheet."""
    try:
        import openpyxl
    except ImportError:
        raise ImportError(
            "XLSX extraction requires openpyxl. "
            "Install with: pip install openpyxl"
        )

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    pages = []
    for page_num, sheet in enumerate(wb.worksheets, 1):
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append("\t".join(cells))
        if rows:
            pages.append((page_num, f"[Sheet: {sheet.title}]\n" + "\n".join(rows)))
    wb.close()
    return pages


def _extract_text(path: Path) -> list[tuple[int, str]]:
    """Extract text from plain text files (TXT, MD, RST, CSV)."""
    text = path.read_text(encoding="utf-8", errors="replace").strip()
    if not text:
        return []
    # Group into virtual pages of ~200 lines
    lines = text.splitlines()
    _PAGE_LINES = 200
    pages = []
    for i in range(0, max(len(lines), 1), _PAGE_LINES):
        chunk = "\n".join(lines[i : i + _PAGE_LINES]).strip()
        if chunk:
            pages.append((len(pages) + 1, chunk))
    return pages


def _extract_pptx(path: Path) -> list[tuple[int, str]]:
    """Extract text from a PPTX file, one virtual page per slide."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "PPTX extraction requires python-pptx. "
            "Install with: pip install python-pptx"
        )

    prs = Presentation(str(path))
    pages = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            pages.append((i, "\n".join(texts)))
    return pages
